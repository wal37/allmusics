from pathlib import Path
import sys


FALLBACK_SITE_PACKAGES = [
    "/Users/wale/Desktop/softtttt/aiapis/apiweb/.venv/lib/python3.13/site-packages",
    "/Users/wale/Desktop/softtttt/foodup/foodweb/.venv/lib/python3.13/site-packages",
]

for site_path in FALLBACK_SITE_PACKAGES:
    if site_path not in sys.path and Path(site_path).exists():
        sys.path.append(site_path)

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path("/Users/wale/Desktop/musically/musicalweb/ALLMUSICS_PITCH_DECK.pptx")
LOGO = "/Users/wale/Desktop/musically/musicalweb/logo.png"


def rgb(value: str) -> RGBColor:
    value = value.replace("#", "")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


BG = rgb("#06040D")
PANEL = rgb("#100C1E")
PAPER = rgb("#F7F2FF")
INK = rgb("#111018")
MUTED_DARK = rgb("#BCB4D2")
MUTED_LIGHT = rgb("#655E74")
PINK = rgb("#FF5FA2")
CORAL = rgb("#FF8E5E")
VIOLET = rgb("#8F7DFF")
BLUE = rgb("#5AC8FF")
LIME = rgb("#D4FF6D")
AMBER = rgb("#FFB55F")
WHITE = rgb("#FFFFFF")
DEEP = rgb("#120D22")
GRID = rgb("#241B38")


SLIDES = [
    {
        "kind": "cover",
        "title": "AllMusics",
        "body": "A music generation platform for artists, composers, film teams, and studios building songs, stems, moods, and sound ideas at production speed.",
        "label": "Cover",
        "dark": True,
        "bars": [PINK, AMBER, VIOLET],
    },
    {
        "kind": "cards",
        "title": "The problem",
        "body": "Music creation tools are still split across ideation, generation, editing, versioning, and delivery. That slows artists, producers, and media teams that need fast output with usable control.",
        "label": "Problem",
        "dark": False,
        "accent": VIOLET,
        "cards": [
            ("Fragmented workflow", "Prompting, arranging, editing, and export often live in different tools."),
            ("Weak production control", "Creators get outputs but not enough structure for real creative iteration."),
            ("Slow commercial delivery", "Teams still need a system for versions, stems, and client-ready assets."),
        ],
    },
    {
        "kind": "cards",
        "title": "Why now",
        "body": "Generative audio is moving from novelty toward production use. The next winner is not only a model, but the platform that wraps real workflow around it.",
        "label": "Why now",
        "dark": True,
        "accent": LIME,
        "cards": [
            ("Creators want speed", "Artists and producers need more throughput without killing momentum."),
            ("Media demand is rising", "Video, branded content, games, and social media need constant audio output."),
            ("Teams need control", "Studios and platforms need repeatability, rights, and export discipline."),
        ],
    },
    {
        "kind": "cards",
        "title": "Product",
        "body": "AllMusics is a prompt-to-sound platform that turns text, references, lyrics, and mood into songs, cues, stems, and editable audio directions from one operating layer.",
        "label": "Platform",
        "dark": False,
        "accent": CORAL,
        "cards": [
            ("Song generation", "Generate complete tracks and musical drafts from prompts."),
            ("Stem control", "Separate output into layers that can be refined and reused."),
            ("Scene scoring", "Create cinematic and ambient pieces for film and visual work."),
            ("Project memory", "Save prompts, versions, and directions across creative sessions."),
        ],
    },
    {
        "kind": "steps",
        "title": "Workflow",
        "body": "The operating rhythm is direction, generation, iteration, and delivery. AllMusics is built around how real audio work moves, not just how a model responds.",
        "label": "Workflow",
        "dark": True,
        "accent": BLUE,
        "steps": [
            "Input style, mood, tempo, references, or lyrics.",
            "Generate candidate outputs and compare directions.",
            "Refine structure, intensity, length, or instrumentation.",
            "Export tracks, stems, and usable commercial assets.",
        ],
    },
    {
        "kind": "cards",
        "title": "Who it serves",
        "body": "AllMusics serves creators and teams that need audio output quickly, but still care about creative quality, asset organization, and delivery speed.",
        "label": "Customers",
        "dark": False,
        "accent": PINK,
        "cards": [
            ("Artists", "Prototype songs, hooks, and demos faster."),
            ("Producers", "Test styles, references, and arrangements."),
            ("Film teams", "Generate soundtrack drafts and scene moods."),
            ("Studios", "Manage version-heavy music workflows."),
        ],
    },
    {
        "kind": "cards",
        "title": "Technology stack",
        "body": "The platform stack includes music generation models, prompt conditioning, vocal and arrangement tooling, version storage, export controls, and future API delivery for external products.",
        "label": "Technology",
        "dark": True,
        "accent": LIME,
        "cards": [
            ("Generation layer", "Audio models for music, vocals, and style variation."),
            ("Editing layer", "Extension, clipping, versioning, and stem operations."),
            ("Business layer", "Workspaces, licensing logic, export, and usage control."),
        ],
    },
    {
        "kind": "cards",
        "title": "Use cases",
        "body": "AllMusics can power direct creator workflows and become embedded inside broader software products that need music generation and audio ideation.",
        "label": "Use cases",
        "dark": False,
        "accent": VIOLET,
        "cards": [
            ("Creator tool", "Direct product for artists and producers."),
            ("Studio engine", "Internal workflow layer for music teams."),
            ("Media pipeline", "Faster soundtrack and branded content generation."),
            ("API platform", "Generation endpoints for partner products."),
        ],
    },
    {
        "kind": "cards",
        "title": "Business model",
        "body": "The platform monetizes through recurring creator subscriptions, team workspaces, commercial licensing, and API usage for software and media partners.",
        "label": "Business",
        "dark": True,
        "accent": AMBER,
        "cards": [
            ("Creator subscriptions", "Starter and Pro plans with monthly credit pools."),
            ("Teams and studios", "Seats, approvals, workspaces, and shared project history."),
            ("Enterprise and API", "Embedded generation, invoicing, and commercial partnerships."),
        ],
    },
    {
        "kind": "cards",
        "title": "Subscription design",
        "body": "The pricing structure is built to convert solo creators first, then expand into teams and platform customers as audio generation becomes operationally important.",
        "label": "Pricing",
        "dark": False,
        "accent": PINK,
        "cards": [
            ("Starter", "$19 per month for idea velocity, drafts, and saved sessions."),
            ("Pro", "$79 per month for heavier generation, stem export, and release workflows."),
            ("Studio", "Custom contracts for seats, licensing controls, reporting, and API access."),
        ],
    },
    {
        "kind": "cards",
        "title": "Differentiation",
        "body": "The defensible layer is workflow depth: generation plus control, versioning, stems, export discipline, and team-ready operations around creative output.",
        "label": "Moat",
        "dark": False,
        "accent": BLUE,
        "cards": [
            ("Not just output", "AllMusics focuses on usable production workflow."),
            ("Not just individuals", "The system is built for teams and commercial work."),
            ("Not just one model", "The platform layer can evolve with the generation stack."),
        ],
    },
    {
        "kind": "steps",
        "title": "Roadmap",
        "body": "Start with generation and creator workflows, then expand into deeper editing, team controls, and platform APIs for external software and media products.",
        "label": "Roadmap",
        "dark": True,
        "accent": PINK,
        "steps": [
            "Phase 1: Creator product and core song generation.",
            "Phase 2: Stem tools, scene scoring, and commercial exports.",
            "Phase 3: Team workspaces, licensing controls, and API access.",
        ],
    },
    {
        "kind": "cover",
        "title": "Vision",
        "body": "AllMusics can become the operating system for fast audio creation: a place where music, soundtrack, ideation, and delivery happen in one modern workflow.",
        "label": "Close",
        "dark": False,
        "bars": [LIME, BLUE, PINK],
    },
]


def add_rect(slide, left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    return shape


def add_round_rect(slide, left, top, width, height, fill, line=None, radius=0.08):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    shape.adjustments[0] = radius
    return shape


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text,
    size=20,
    color=INK,
    bold=False,
    font="Aptos",
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.vertical_anchor = valign
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = text
    run = p.runs[0]
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_badge(slide, left, top, text, fill, color):
    width = max(1.7, 0.11 * len(text) + 0.7)
    add_round_rect(slide, left, top, Inches(width), Inches(0.42), fill, None, 0.18)
    add_text(
        slide,
        left + Inches(0.18),
        top + Inches(0.08),
        Inches(width - 0.25),
        Inches(0.2),
        text.upper(),
        size=10,
        color=color,
        bold=True,
        font="Aptos Display",
    )


def add_frame(slide, num, dark):
    base = BG if dark else PAPER
    border = GRID if dark else rgb("#DDD4EF")
    add_rect(slide, 0, 0, Inches(13.333), Inches(7.5), base)
    add_rect(slide, Inches(0.42), Inches(0.38), Inches(12.48), Inches(6.74), base, border)
    add_rect(slide, Inches(0.84), Inches(0.84), Inches(11.65), Inches(0.06), PINK if dark else VIOLET)
    add_text(
        slide,
        Inches(12.0),
        Inches(6.76),
        Inches(0.55),
        Inches(0.18),
        f"{num:02d}",
        size=10,
        color=MUTED_DARK if dark else MUTED_LIGHT,
        bold=True,
        font="Aptos Display",
        align=PP_ALIGN.RIGHT,
    )


def add_logo(slide, left, top, height=0.56):
    if Path(LOGO).exists():
        slide.shapes.add_picture(LOGO, left, top, height=Inches(height))


def add_card(slide, left, top, width, height, title, body, accent, dark):
    fill = PANEL if dark else WHITE
    line = None if dark else rgb("#DDD4EF")
    body_color = MUTED_DARK if dark else MUTED_LIGHT
    title_color = WHITE if dark else INK
    add_round_rect(slide, left, top, width, height, fill, line, 0.1)
    add_rect(slide, left, top, width, Inches(0.06), accent)
    add_text(slide, left + Inches(0.18), top + Inches(0.18), width - Inches(0.3), Inches(0.34), title, size=16, color=title_color, bold=True, font="Aptos Display")
    add_text(slide, left + Inches(0.18), top + Inches(0.58), width - Inches(0.3), height - Inches(0.74), body, size=11.5, color=body_color)


def add_step(slide, left, top, index, text, accent, dark):
    text_color = WHITE if dark else INK
    body_color = MUTED_DARK if dark else MUTED_LIGHT
    add_round_rect(slide, left, top, Inches(0.42), Inches(0.42), accent, None, 0.14)
    add_text(slide, left, top + Inches(0.08), Inches(0.42), Inches(0.16), f"{index:02d}", size=11, color=INK, bold=True, font="Aptos Display", align=PP_ALIGN.CENTER)
    add_text(slide, left, top + Inches(0.56), Inches(2.25), Inches(0.95), text, size=13, color=body_color)
    add_text(slide, left + Inches(0.62), top + Inches(0.02), Inches(1.6), Inches(0.25), "Step", size=10, color=accent, bold=True, font="Aptos Display")
    add_text(slide, left, top + Inches(0.56), Inches(2.25), Inches(0.95), text, size=13, color=text_color if False else body_color)


def draw_cover(slide, data, num):
    dark = data["dark"]
    add_frame(slide, num, dark)
    add_logo(slide, Inches(1.08), Inches(1.0), 0.72)
    add_badge(slide, Inches(1.08), Inches(1.88), data["label"], LIME if dark else PINK, INK if dark else WHITE)
    title_color = WHITE if dark else INK
    body_color = MUTED_DARK if dark else MUTED_LIGHT
    add_text(slide, Inches(1.08), Inches(2.48), Inches(5.6), Inches(0.7), data["title"], size=31, color=title_color, bold=True, font="Aptos Display")
    add_text(slide, Inches(1.08), Inches(3.28), Inches(5.8), Inches(1.15), data["body"], size=18, color=body_color)
    x = Inches(8.0)
    widths = [Inches(3.25), Inches(2.45), Inches(3.0)]
    heights = [Inches(0.82), Inches(0.7), Inches(2.05)]
    tops = [Inches(1.22), Inches(2.38), Inches(3.34)]
    offsets = [Inches(0), Inches(0.28), Inches(0)]
    for idx, color in enumerate(data["bars"]):
        add_round_rect(slide, x + offsets[idx], tops[idx], widths[idx], heights[idx], color, None, 0.06)


def draw_cards(slide, data, num):
    dark = data["dark"]
    add_frame(slide, num, dark)
    title_color = WHITE if dark else INK
    body_color = MUTED_DARK if dark else MUTED_LIGHT
    add_badge(slide, Inches(1.08), Inches(1.0), data["label"], data["accent"], INK)
    add_text(slide, Inches(1.08), Inches(1.56), Inches(6.4), Inches(0.78), data["title"], size=27, color=title_color, bold=True, font="Aptos Display")
    add_text(slide, Inches(1.08), Inches(2.38), Inches(6.5), Inches(1.0), data["body"], size=14.5, color=body_color)

    cards = data["cards"]
    if len(cards) == 4:
        positions = [
            (Inches(1.08), Inches(3.48), Inches(2.85), Inches(1.7)),
            (Inches(4.12), Inches(3.48), Inches(2.85), Inches(1.7)),
            (Inches(7.16), Inches(3.48), Inches(2.85), Inches(1.7)),
            (Inches(10.2), Inches(3.48), Inches(2.05), Inches(1.7)),
        ]
    else:
        positions = [
            (Inches(1.08), Inches(3.58), Inches(3.45), Inches(1.72)),
            (Inches(4.86), Inches(3.58), Inches(3.45), Inches(1.72)),
            (Inches(8.64), Inches(3.58), Inches(3.45), Inches(1.72)),
        ]

    for (title, body), (left, top, width, height) in zip(cards, positions):
        add_card(slide, left, top, width, height, title, body, data["accent"], dark)


def draw_steps(slide, data, num):
    dark = data["dark"]
    add_frame(slide, num, dark)
    title_color = WHITE if dark else INK
    body_color = MUTED_DARK if dark else MUTED_LIGHT
    add_badge(slide, Inches(1.08), Inches(1.0), data["label"], data["accent"], INK)
    add_text(slide, Inches(1.08), Inches(1.56), Inches(6.2), Inches(0.78), data["title"], size=27, color=title_color, bold=True, font="Aptos Display")
    add_text(slide, Inches(1.08), Inches(2.38), Inches(6.7), Inches(1.0), data["body"], size=14.5, color=body_color)

    x_positions = [Inches(1.3), Inches(4.02), Inches(6.74), Inches(9.46)]
    for idx, (step, left) in enumerate(zip(data["steps"], x_positions), start=1):
        add_round_rect(slide, left, Inches(3.86), Inches(2.28), Inches(1.36), PANEL if dark else WHITE, None if dark else rgb("#DDD4EF"), 0.08)
        add_rect(slide, left, Inches(3.86), Inches(0.1), Inches(1.36), data["accent"])
        add_round_rect(slide, left + Inches(0.18), Inches(4.1), Inches(0.42), Inches(0.42), data["accent"], None, 0.14)
        add_text(slide, left + Inches(0.01), Inches(4.18), Inches(0.4), Inches(0.14), f"{idx:02d}", size=11, color=INK, bold=True, font="Aptos Display", align=PP_ALIGN.CENTER)
        add_text(slide, left + Inches(0.18), Inches(4.66), Inches(1.84), Inches(0.72), step, size=12.5, color=body_color)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6]

    for idx, data in enumerate(SLIDES, start=1):
        slide = prs.slides.add_slide(layout)
        if data["kind"] == "cover":
            draw_cover(slide, data, idx)
        elif data["kind"] == "cards":
            draw_cards(slide, data, idx)
        elif data["kind"] == "steps":
            draw_steps(slide, data, idx)

    prs.core_properties.title = "AllMusics Pitch Deck"
    prs.core_properties.author = "OpenAI"
    prs.core_properties.last_modified_by = "OpenAI"
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
